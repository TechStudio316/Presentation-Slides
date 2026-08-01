import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# -----------------------------------------------------------------------------
# CONSTANTS & BRANDING SETUP
# -----------------------------------------------------------------------------
COLOR_ROYAL_BLUE = RGBColor(10, 78, 163)   # #0A4EA3
COLOR_GOLD       = RGBColor(245, 179, 1)   # #F5B301
COLOR_WHITE      = RGBColor(255, 255, 255)
COLOR_DARK_TEXT  = RGBColor(30, 41, 59)    # Slate 800
COLOR_LIGHT_BG   = RGBColor(248, 250, 252)  # Slate 50
COLOR_CARD_BORDER= RGBColor(226, 232, 240) # Slate 200
COLOR_MUTED_TEXT = RGBColor(100, 116, 139) # Slate 500


FONT_HEADING = "Segoe UI"
FONT_BODY    = "Segoe UI"


BRAND_NAME   = "BLUESKY EDUCATIONAL CONSULTS"
BRAND_TAGLINE= "Learning Beyond Borders"
FOOTER_TEXT  = f"{BRAND_NAME} | {BRAND_TAGLINE}"


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[6]


# -----------------------------------------------------------------------------
# HELPER FUNCTIONS
# -----------------------------------------------------------------------------
def set_bg_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header(slide, title_text, category_text="YEAR 7 ENGLISH • KEY STAGE 3"):
    # Header Category Tag
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = FONT_HEADING
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_GOLD


    # Main Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_ROYAL_BLUE


def add_footer(slide):
    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_CARD_BORDER
    line.line.fill.background()


    # Footer Text
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.95), Inches(11.733), Inches(0.4))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = FOOTER_TEXT
    p.font.name = FONT_BODY
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_MUTED_TEXT


def create_card(slide, left, top, width, height, bg_color=COLOR_WHITE, border_color=COLOR_CARD_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text


# -----------------------------------------------------------------------------
# SLIDE BUILDERS
# -----------------------------------------------------------------------------


# SLIDE 1: Cover Slide
def build_slide_1():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_ROYAL_BLUE)


    # Decorative Accent Bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.15), Inches(4.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_GOLD
    accent.line.fill.background()


    # Title & Subtitle Box
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.5), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True


    p0 = tf.paragraphs[0]
    p0.text = BRAND_NAME
    p0.font.name = FONT_HEADING
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_GOLD


    p1 = tf.add_paragraph()
    p1.text = "Introduction to English Skills"
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_before = Pt(20)


    p2 = tf.add_paragraph()
    p2.text = "Understanding Language, Literature & Effective Communication"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(226, 232, 240)
    p2.space_before = Pt(10)


    p3 = tf.add_paragraph()
    p3.text = "Year 7 English • Key Stage 3 • Week 1, Lesson 1 (90 Minutes)"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_GOLD
    p3.space_before = Pt(40)


    add_notes(slide, "TEACHER NOTES: Welcome students warmly. Introduce Bluesky Educational Consults. Set expectations for interactive participation.")


# SLIDE 2: Welcome Back
def build_slide_2():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Welcome to Key Stage 3 English!")
    add_footer(slide)


    # Left Box: Welcome Message
    create_card(slide, Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.0))
    tb = slide.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(6.9), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True


    p = tf.paragraphs[0]
    p.text = "A New Chapter Begins"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_ROYAL_BLUE


    bullets = [
        "Welcome to your Secondary English journey!",
        "In Key Stage 3, you will explore powerful stories, debate big ideas, and master the art of expression.",
        "Our classroom is a safe space to test ideas, make mistakes, and grow as confident communicators.",
        "Get ready to read critically, write creatively, and speak persuasively."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_DARK_TEXT
        p.space_before = Pt(12)


    # Right Box: Expectations Card
    create_card(slide, Inches(8.6), Inches(1.6), Inches(3.9), Inches(5.0), bg_color=COLOR_ROYAL_BLUE)
    tb2 = slide.shapes.add_textbox(Inches(8.8), Inches(1.8), Inches(3.5), Inches(4.6))
    tf2 = tb2.text_frame
    tf2.word_wrap = True


    p2 = tf2.paragraphs[0]
    p2.text = "Classroom Values"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_GOLD


    vals = ["📖 Active Listening", "🙋 Active Participation", "🤝 Mutual Respect", "💡 Creative Curiosity", "✏️ High Endeavour"]
    for v in vals:
        p = tf2.add_paragraph()
        p.text = v
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(16)


    add_notes(slide, "TEACHER NOTES (5 mins): Warm welcome. Go over learning culture. Emphasise that wrong answers are stepping stones to learning.")


# SLIDE 3: Ice Breaker
def build_slide_3():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Ice Breaker: What Makes a Good Communicator?")
    add_footer(slide)


    # Main Question Card
    create_card(slide, Inches(0.8), Inches(1.6), Inches(11.733), Inches(1.8), bg_color=COLOR_ROYAL_BLUE)
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💬 Think of someone you love listening to (a YouTuber, teacher, author, or friend).\nWhat makes them so effective when they communicate?"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE


    # 3 Discussion Prompts
    prompts = [
        ("1. Tone & Voice", "Is it how they sound, the words they choose, or their body language?"),
        ("2. Clarity", "Do they make complex ideas easy to understand?"),
        ("3. Engagement", "How do they keep your attention from start to finish?")
    ]
    for i, (title, desc) in enumerate(prompts):
        left = Inches(0.8 + i * 4.0)
        create_card(slide, left, Inches(3.6), Inches(3.733), Inches(3.0))
        tb_p = slide.shapes.add_textbox(left + Inches(0.2), Inches(3.8), Inches(3.333), Inches(2.6))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True


        p_t = tf_p.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_ROYAL_BLUE


        p_d = tf_p.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = COLOR_DARK_TEXT
        p_d.space_before = Pt(10)


    add_notes(slide, "TEACHER NOTES (8 mins): Ask 3-4 students to share their ideas. Write key words (e.g. eye contact, vocabulary, confidence) on the board.")


# SLIDE 4: Lesson Objectives
def build_slide_4():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Learning Objectives")
    add_footer(slide)


    objectives = [
        ("01", "Distinguish Core Domains", "Explain clearly what English Language and English Literature are and how they connect."),
        ("02", "Identify the Four Pillars", "Recognise and apply the four key language skills: Reading, Writing, Speaking, and Listening."),
        ("03", "Understand Communication", "Recognise why effective communication is vital in education, careers, and everyday life."),
        ("04", "Master Key Vocabulary", "Use subject-specific terminology (Author, Audience, Purpose, Tone, Genre) with confidence."),
        ("05", "Analyse Text for Meaning", "Read a short literary passage and extract both literal and inferential meaning using evidence.")
    ]


    for i, (num, title, desc) in enumerate(objectives):
        top = Inches(1.6 + i * 1.0)
        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(0.8), Inches(0.8))
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_ROYAL_BLUE
        badge.line.fill.background()
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = num
        p_b.font.size = Pt(14)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_GOLD
        p_b.alignment = PP_ALIGN.CENTER


        # Content Card
        create_card(slide, Inches(1.8), top, Inches(10.733), Inches(0.8))
        tb = slide.shapes.add_textbox(Inches(2.0), top + Inches(0.05), Inches(10.3), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True


        p_t = tf.paragraphs[0]
        p_t.text = title + " — "
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_ROYAL_BLUE


        run = p_t.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = COLOR_DARK_TEXT


    add_notes(slide, "TEACHER NOTES (3 mins): Read through the objectives. Emphasise that today builds the foundation for all Key Stage 3 English.")


# SLIDE 5: Success Criteria
def build_slide_5():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Success Criteria")
    add_footer(slide)


    criteria = [
        ("Bronze Level", "I can state the difference between English Language and English Literature, and list the four skills.", COLOR_ROYAL_BLUE),
        ("Silver Level", "I can define key vocabulary (Audience, Purpose, Tone) and identify them in a short text.", COLOR_ROYAL_BLUE),
        ("Gold Level", "I can explain the author's purpose using explicit evidence and make simple inferences from text.", COLOR_GOLD)
    ]


    for i, (level, desc, color) in enumerate(criteria):
        top = Inches(1.6 + i * 1.7)
        create_card(slide, Inches(0.8), top, Inches(11.733), Inches(1.4))


        # Color Accent Line on Left of Card
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(0.2), Inches(1.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()


        tb = slide.shapes.add_textbox(Inches(1.2), top + Inches(0.15), Inches(11.1), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True


        p0 = tf.paragraphs[0]
        p0.text = f"✓ {level}"
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = color


        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(14)
        p1.font.color.rgb = COLOR_DARK_TEXT
        p1.space_before = Pt(6)


    add_notes(slide, "TEACHER NOTES (2 mins): Explain that students should aim for Gold today. Self-assessment will happen at the end.")


# SLIDE 6: Warm-up Activity
def build_slide_6():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Warm-up: Word Association Challenge", "STARTER ACTIVITY • 5 MINS")
    add_footer(slide)


    # Instruction Banner
    create_card(slide, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.9), bg_color=COLOR_ROYAL_BLUE)
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "⚡ Challenge: You have 2 minutes! Write down as many words as you associate with each pillar below:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE


    pillars = [
        ("📖 READING", ["Books", "Decipher", "Analyze", "Comprehension", "Library", "Context"]),
        ("✍️ WRITING", ["Grammar", "Punctuation", "Essays", "Creativity", "Vocabulary", "Draft"]),
        ("🗣️ SPEAKING", ["Debate", "Presentation", "Tone", "Pitch", "Expression", "Dialogue"]),
        ("👂 LISTENING", ["Attention", "Empathy", "Understanding", "Notes", "Active", "Focus"])
    ]


    for i, (title, examples) in enumerate(pillars):
        left = Inches(0.8 + i * 3.0)
        create_card(slide, left, Inches(2.6), Inches(2.733), Inches(4.0))


        tb_c = slide.shapes.add_textbox(left + Inches(0.15), Inches(2.8), Inches(2.433), Inches(3.6))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True


        p_t = tf_c.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_ROYAL_BLUE


        p_sub = tf_c.add_paragraph()
        p_sub.text = "Examples:"
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_GOLD
        p_sub.space_before = Pt(10)


        for ex in examples:
            p_ex = tf_c.add_paragraph()
            p_ex.text = f"• {ex}"
            p_ex.font.size = Pt(12)
            p_ex.font.color.rgb = COLOR_DARK_TEXT
            p_ex.space_before = Pt(4)


    add_notes(slide, "TEACHER NOTES (5 mins): Set a timer for 2 mins. Rapid fire answers from students. Reveal examples on slide.")


# SLIDE 7: What is English?
def build_slide_7():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "What is English? Language vs. Literature")
    add_footer(slide)


    # Column 1: Language
    create_card(slide, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.0))
    tb1 = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.6))
    tf1 = tb1.text_frame
    tf1.word_wrap = True


    p1 = tf1.paragraphs[0]
    p1.text = "🔤 English Language"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_ROYAL_BLUE


    p1_sub = tf1.add_paragraph()
    p1_sub.text = "The Tool of Communication"
    p1_sub.font.size = Pt(12)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = COLOR_GOLD
    p1_sub.space_before = Pt(4)


    lang_points = [
        "Focuses on HOW words work together.",
        "Includes grammar, punctuation, vocabulary, and sentence structures.",
        "Examines how language changes according to audience and purpose.",
        "Teaches non-fiction writing: articles, speeches, letters, and reports."
    ]
    for pt in lang_points:
        p = tf1.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK_TEXT
        p.space_before = Pt(10)


    # Column 2: Literature
    create_card(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.0))
    tb2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.6))
    tf2 = tb2.text_frame
    tf2.word_wrap = True


    p2 = tf2.paragraphs[0]
    p2.text = "📚 English Literature"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_ROYAL_BLUE


    p2_sub = tf2.add_paragraph()
    p2_sub.text = "The Study of Human Experience"
    p2_sub.font.size = Pt(12)
    p2_sub.font.bold = True
    p2_sub.font.color.rgb = COLOR_GOLD
    p2_sub.space_before = Pt(4)


    lit_points = [
        "Focuses on WHAT writers express and WHY.",
        "Includes novels, plays, poetry, and historical texts.",
        "Explores themes like love, power, conflict, and identity.",
        "Teaches analytical thinking and critical interpretation."
    ]
    for pt in lit_points:
        p = tf2.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK_TEXT
        p.space_before = Pt(10)


    add_notes(slide, "TEACHER NOTES (7 mins): Highlight that both sides depend on each other. You need Language mechanics to write and interpret Literature.")


# SLIDE 8: The Four Language Skills
def build_slide_8():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "The Four Language Skills")
    add_footer(slide)


    skills = [
        ("📖 READING", "Receptive Skill", "Decoding symbols, interpreting meaning, analyzing techniques, and building comprehension.", COLOR_ROYAL_BLUE),
        ("✍️ WRITING", "Productive Skill", "Structuring thoughts, selecting vocabulary, applying grammar, and crafting narratives.", COLOR_ROYAL_BLUE),
        ("🗣️ SPEAKING", "Productive Skill", "Articulating ideas, adapting tone, presenting arguments, and performing drama.", COLOR_GOLD),
        ("👂 LISTENING", "Receptive Skill", "Processing spoken words, active listening, evaluating speeches, and understanding intent.", COLOR_GOLD)
    ]


    for i, (title, type_label, desc, color) in enumerate(skills):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 6.0)
        top = Inches(1.6 + row * 2.6)


        create_card(slide, left, top, Inches(5.733), Inches(2.3))


        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.333), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True


        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(18)
        p0.font.bold = True
        p0.font.color.rgb = color


        p_sub = tf.add_paragraph()
        p_sub.text = f"[{type_label}]"
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_MUTED_TEXT


        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(13)
        p1.font.color.rgb = COLOR_DARK_TEXT
        p1.space_before = Pt(8)


    add_notes(slide, "TEACHER NOTES (5 mins): Explain Receptive (taking in information) vs Productive (producing information) skills.")


# SLIDE 9: Why English Matters
def build_slide_9():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Why English Matters Beyond the Classroom")
    add_footer(slide)


    areas = [
        ("🎓 Education", "Mastering English unlocks all other subjects—from writing science reports to analyzing history source documents."),
        ("💼 Careers", "Employers consistently rank clear written and verbal communication as the #1 essential workplace skill."),
        ("🌐 Media Literacy", "Helps you navigate social media, spot fake news, and evaluate bias in public information."),
        ("🗣️ Everyday Life", "Expressing your feelings clearly, resolving conflicts, and forming deep relationships require strong communication.")
    ]


    for i, (title, desc) in enumerate(areas):
        left = Inches(0.8 + i * 2.95)
        create_card(slide, left, Inches(1.6), Inches(2.8), Inches(5.0))


        tb = slide.shapes.add_textbox(left + Inches(0.15), Inches(1.8), Inches(2.5), Inches(4.6))
        tf = tb.text_frame
        tf.word_wrap = True


        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_ROYAL_BLUE


        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_DARK_TEXT
        p1.space_before = Pt(12)


    add_notes(slide, "TEACHER NOTES (5 mins): Emphasise real-world relevance. Ask students: 'What job does NOT require communication?' (Spoiler: none!)")


# SLIDE 10: Building Vocabulary
def build_slide_10():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Building Vocabulary: Core English Toolkit", "SUBJECT-SPECIFIC VOCABULARY")
    add_footer(slide)


    terms = [
        ("Author", "The writer of a text, story, or document."),
        ("Audience", "The intended reader, listener, or viewer."),
        ("Purpose", "The reason a text was created (e.g. inform, persuade, entertain)."),
        ("Genre", "The category or style of literature/media (e.g. Sci-Fi, Gothic, Poetry)."),
        ("Tone", "The mood or attitude expressed by the writer's words (e.g. mysterious, joyful).")
    ]


    for i, (term, defn) in enumerate(terms):
        top = Inches(1.6 + i * 1.0)
        # Term Header Card
        create_card(slide, Inches(0.8), top, Inches(2.5), Inches(0.8), bg_color=COLOR_ROYAL_BLUE)
        tb_t = slide.shapes.add_textbox(Inches(0.9), top + Inches(0.15), Inches(2.3), Inches(0.5))
        p_t = tb_t.text_frame.paragraphs[0]
        p_t.text = term
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_GOLD


        # Definition Card
        create_card(slide, Inches(3.4), top, Inches(9.133), Inches(0.8))
        tb_d = slide.shapes.add_textbox(Inches(3.6), top + Inches(0.15), Inches(8.7), Inches(0.5))
        p_d = tb_d.text_frame.paragraphs[0]
        p_d.text = defn
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = COLOR_DARK_TEXT


    add_notes(slide, "TEACHER NOTES (6 mins): Have students copy these definitions into their exercise books. These 5 words will be used every lesson.")


# SLIDE 11: Reading Activity
def build_slide_11():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Reading Activity: The Discovery", "TEXT ANALYSIS • 10 MINS")
    add_footer(slide)


    # Passage Card
    create_card(slide, Inches(0.8), Inches(1.6), Inches(11.733), Inches(4.0))


    tb = slide.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.133), Inches(3.6))
    tf = tb.text_frame
    tf.word_wrap = True


    p_head = tf.paragraphs[0]
    p_head.text = "Excerpt from 'The Forgotten Library' by Maya Lin"
    p_head.font.size = Pt(12)
    p_head.font.italic = True
    p_head.font.color.rgb = COLOR_MUTED_TEXT


    passage = (
        "The heavy wooden door groaned as Arthur pushed it open, revealing a cavernous room "
        "bathed in golden dust motes. Bookshelves reached up toward vaulted ceilings like ancient oak trees, "
        "their branches heavy with leather-bound volumes. A gentle smell of old paper and dried lavender "
        "hung in the still air. In the centre of the room stood a polished mahogany desk, and upon it lay "
        "a single brass key resting on an unopened letter. Arthur took a slow step forward, his heart "
        "hammering against his ribs like a trapped bird."
    )


    p_body = tf.add_paragraph()
    p_body.text = passage
    p_body.font.size = Pt(16)
    p_body.font.color.rgb = COLOR_DARK_TEXT
    p_body.space_before = Pt(14)


    # Callout Banner below passage
    create_card(slide, Inches(0.8), Inches(5.8), Inches(11.733), Inches(0.8), bg_color=COLOR_ROYAL_BLUE)
    tb_c = slide.shapes.add_textbox(Inches(1.0), Inches(5.9), Inches(11.333), Inches(0.6))
    p_c = tb_c.text_frame.paragraphs[0]
    p_c.text = "🔍 Task: Read the passage silently twice. Notice how the author creates atmosphere using sensory details."
    p_c.font.size = Pt(13)
    p_c.font.bold = True
    p_c.font.color.rgb = COLOR_WHITE


    add_notes(slide, "TEACHER NOTES (5 mins): Read aloud once for fluency and expression. Then ask a confident student to read it.")


# SLIDE 12: Teacher Modelling
def build_slide_12():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Teacher Modelling: Think Aloud Strategy", "I DO • ANNOTATING TEXT")
    add_footer(slide)


    # Split: Left side = Text snippet, Right side = Thinking Process
    create_card(slide, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.0))
    tb_l = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.6))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True


    p0 = tf_l.paragraphs[0]
    p0.text = "Text Focus:"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_GOLD


    p1 = tf_l.add_paragraph()
    p1.text = '"...his heart hammering against his ribs like a trapped bird."'
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_ROYAL_BLUE
    p1.space_before = Pt(10)


    create_card(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.0), bg_color=COLOR_ROYAL_BLUE)
    tb_r = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.6))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True


    p2 = tf_r.paragraphs[0]
    p2.text = "🧠 Teacher's 'Think Aloud':"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_GOLD


    thoughts = [
        "1. Identify Technique: 'like a trapped bird' is a Simile.",
        "2. Analyze Imagery: A trapped bird flutters frantically and wants to escape.",
        "3. Connect to Meaning: Arthur isn't just nervous; he feels overwhelmed and trapped by mystery.",
        "4. Evaluate Atmosphere: The writer creates tension and suspense."
    ]
    for th in thoughts:
        p = tf_r.add_paragraph()
        p.text = th
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(12)


    add_notes(slide, "TEACHER NOTES (5 mins): Model how good readers think while reading. Demonstrate annotating on a whiteboard.")


# SLIDE 13: Guided Reading
def build_slide_13():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Guided Reading: Let's Unpack Together", "WE DO • GUIDED ANALYSIS")
    add_footer(slide)


    questions = [
        ("Question 1: Sensory Language", "Which two senses (sight, sound, smell, touch, taste) does the author appeal to in line 3?", "Smell ('old paper and dried lavender') and Sound ('door groaned')."),
        ("Question 2: Figurative Language", "What metaphor is used to describe the bookshelves in the room?", "They are compared to 'ancient oak trees' with branches heavy with books."),
        ("Question 3: Mystery & Suspense", "Which object in the room signals to the reader that a story is about to begin?", "The brass key resting on the unopened letter.")
    ]


    for i, (q_title, q_text, q_ans) in enumerate(questions):
        top = Inches(1.6 + i * 1.7)
        create_card(slide, Inches(0.8), top, Inches(11.733), Inches(1.5))


        tb = slide.shapes.add_textbox(Inches(1.0), top + Inches(0.1), Inches(11.333), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True


        p0 = tf.paragraphs[0]
        p0.text = q_title
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_ROYAL_BLUE


        p1 = tf.add_paragraph()
        p1.text = q_text
        p1.font.size = Pt(13)
        p1.font.color.rgb = COLOR_DARK_TEXT


        p2 = tf.add_paragraph()
        p2.text = f"💡 Answer: {q_ans}"
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(16, 185, 129) # Emerald Green for answer
        p2.space_before = Pt(4)


    add_notes(slide, "TEACHER NOTES (8 mins): Solicit answers from volunteers. Click to reveal/highlight answers on the slide.")


# SLIDE 14: Comprehension Questions
def build_slide_14():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Comprehension Challenge: Literal vs. Inferential", "YOU DO • CHECKING UNDERSTANDING")
    add_footer(slide)


    # Left Column: Literal Questions
    create_card(slide, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.0))
    tb_l = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.3), Inches(4.6))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True


    p_l_head = tf_l.paragraphs[0]
    p_l_head.text = "📍 Literal Questions (In the text)"
    p_l_head.font.size = Pt(16)
    p_l_head.font.bold = True
    p_l_head.font.color.rgb = COLOR_ROYAL_BLUE


    l_qs = [
        "1. What kind of door did Arthur push open?\n   → A heavy wooden door.",
        "2. What object was sitting on top of the mahogany desk?\n   → A single brass key and an unopened letter."
    ]
    for q in l_qs:
        p = tf_l.add_paragraph()
        p.text = q
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK_TEXT
        p.space_before = Pt(14)


    # Right Column: Inferential Questions
    create_card(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.0))
    tb_r = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.3), Inches(4.6))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True


    p_r_head = tf_r.paragraphs[0]
    p_r_head.text = "🔍 Inferential Questions (Read between lines)"
    p_r_head.font.size = Pt(16)
    p_r_head.font.bold = True
    p_r_head.font.color.rgb = COLOR_ROYAL_BLUE


    r_qs = [
        "3. Has anyone been in this library recently? How do you know?\n   → No. The presence of 'golden dust motes' and 'still air' suggests it was untouched.",
        "4. How does Arthur feel about being in the library?\n   → He feels both awed and terrified ('heart hammering')."
    ]
    for q in r_qs:
        p = tf_r.add_paragraph()
        p.text = q
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK_TEXT
        p.space_before = Pt(14)


    add_notes(slide, "TEACHER NOTES (8 mins): Explain difference: Literal = right there on page. Inferential = reader detective work using clues.")


# SLIDE 15: Interactive Vocabulary Activity
def build_slide_15():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Vocabulary Matching Challenge", "INTERACTIVE ACTIVITY")
    add_footer(slide)


    # Left Box: Terms
    create_card(slide, Inches(0.8), Inches(1.6), Inches(4.0), Inches(5.0), bg_color=COLOR_ROYAL_BLUE)
    tb_l = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(3.6), Inches(4.6))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.paragraphs[0].text = "TERMS"
    tf_l.paragraphs[0].font.size = Pt(16)
    tf_l.paragraphs[0].font.bold = True
    tf_l.paragraphs[0].font.color.rgb = COLOR_GOLD


    terms_list = ["1. Cavernous", "2. Vaulted", "3. Hammering", "4. Stillness"]
    for t in terms_list:
        p = tf_l.add_paragraph()
        p.text = t
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(20)


    # Right Box: Definitions (Shuffled / Matched)
    create_card(slide, Inches(5.1), Inches(1.6), Inches(7.433), Inches(5.0))
    tb_r = slide.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(7.0), Inches(4.6))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.paragraphs[0].text = "DEFINITIONS & MATCHES"
    tf_r.paragraphs[0].font.size = Pt(16)
    tf_r.paragraphs[0].font.bold = True
    tf_r.paragraphs[0].font.color.rgb = COLOR_ROYAL_BLUE


    defs_list = [
        "A. Beating rapidly and loudly against something.  ➔ [Matches #3]",
        "B. Vast, huge, and dark like a deep cave.  ➔ [Matches #1]",
        "C. Quiet, calm, and undisturbed silence.  ➔ [Matches #4]",
        "D. Constructed with high, arched ceilings or roofs.  ➔ [Matches #2]"
    ]
    for d in defs_list:
        p = tf_r.add_paragraph()
        p.text = d
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK_TEXT
        p.space_before = Pt(16)


    add_notes(slide, "TEACHER NOTES (5 mins): Have students call out the matching numbers and letters before revealing final matches.")


# SLIDE 16: Think–Pair–Share
def build_slide_16():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Think–Pair–Share: Author's Purpose", "COLLABORATIVE TASK • 6 MINS")
    add_footer(slide)


    steps = [
        ("THINK (1 Min)", "Consider individually:\nWhy did the author, Maya Lin, write this opening passage? What emotion did she want you to feel?", COLOR_ROYAL_BLUE),
        ("PAIR (2 Mins)", "Turn to your partner:\nCompare your ideas. Did you pick out the same descriptive words? Agree on the main purpose.", COLOR_ROYAL_BLUE),
        ("SHARE (3 Mins)", "Class Discussion:\nBe ready to share your partner's best idea with the class using evidence from the text.", COLOR_GOLD)
    ]


    for i, (title, desc, color) in enumerate(steps):
        left = Inches(0.8 + i * 4.0)
        create_card(slide, left, Inches(1.6), Inches(3.733), Inches(5.0))


        # Top Accent Header
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.6), Inches(3.733), Inches(0.8))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        tf_h = hdr.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = title
        p_h.font.size = Pt(14)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.alignment = PP_ALIGN.CENTER


        tb = slide.shapes.add_textbox(left + Inches(0.2), Inches(2.6), Inches(3.333), Inches(3.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK_TEXT


    add_notes(slide, "TEACHER NOTES (6 mins): Enforce strict timings with a visual timer on screen. Circulate room during Pair phase.")


# SLIDE 17: Independent Practice
def build_slide_17():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Independent Practice: Written Analysis", "YOU DO • INDIVIDUAL WORK • 10 MINS")
    add_footer(slide)


    # Task Box
    create_card(slide, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0))
    tb = slide.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.133), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True


    p0 = tf.paragraphs[0]
    p0.text = "📝 Task: Answer the following prompt in full sentences in your exercise notebook:"
    p0.font.size = Pt(15)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ROYAL_BLUE


    p_prompt = tf.add_paragraph()
    p_prompt.text = '"How does Maya Lin create a sense of mystery and wonder in \'The Forgotten Library\'?"'
    p_prompt.font.size = Pt(16)
    p_prompt.font.bold = True
    p_prompt.font.color.rgb = COLOR_GOLD
    p_prompt.space_before = Pt(10)


    p_guide = tf.add_paragraph()
    p_guide.text = "Requirements for your paragraph:"
    p_guide.font.size = Pt(13)
    p_guide.font.bold = True
    p_guide.font.color.rgb = COLOR_DARK_TEXT
    p_guide.space_before = Pt(14)


    reqs = [
        "Write at least 3-4 full sentences.",
        "Include at least ONE direct quotation from the text (e.g. 'cavernous room' or 'ancient oak trees').",
        "Use key vocabulary words (e.g. Tone, Imagery, Audience, Atmosphere).",
        "Check your capital letters, full stops, and spelling before submitting."
    ]
    for r in reqs:
        p = tf.add_paragraph()
        p.text = f"✓ {r}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_DARK_TEXT
        p.space_before = Pt(6)


    add_notes(slide, "TEACHER NOTES (10 mins): Silent writing time. Walk around room to offer support to lower-attaining students.")


# SLIDE 18: Creative Challenge
def build_slide_18():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Creative Writing Challenge: Self-Introduction", "CREATIVE APPLICATION • 8 MINS")
    add_footer(slide)


    create_card(slide, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0))
    tb = slide.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.133), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True


    p0 = tf.paragraphs[0]
    p0.text = "🎨 Task: Introduce Yourself using Descriptive Language"
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ROYAL_BLUE


    p1 = tf.add_paragraph()
    p1.text = "Write a short 50-word paragraph introducing yourself to the class, but with a twist:"
    p1.font.size = Pt(14)
    p1.font.color.rgb = COLOR_DARK_TEXT
    p1.space_before = Pt(8)


    rules = [
        "Do NOT just list facts (e.g. 'My name is Alex and I like football').",
        "Use vivid verbs and sensory descriptions (e.g. 'I am a storm of energy on the football pitch...').",
        "Include ONE simile or metaphor to describe your personality.",
        "Set an enthusiastic or mysterious TONE."
    ]
    for r in rules:
        p = tf.add_paragraph()
        p.text = f"⭐ {r}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK_TEXT
        p.space_before = Pt(10)


    add_notes(slide, "TEACHER NOTES (8 mins): Fun activity to gauge initial writing skills. Invite 2-3 students to read theirs aloud.")


# SLIDE 19: Common Mistakes
def build_slide_19():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Common Pitfalls in Year 7 English", "MISCONCEPTIONS TO AVOID")
    add_footer(slide)


    pitfalls = [
        ("❌ Pitfall 1: One-Word Answers", "Giving 'Yes' or 'It was good' answers.", "✔️ Fix: Always answer in full sentences and give reasons for your point."),
        ("❌ Pitfall 2: Floating Quotations", "Dropping quotes into paragraphs without context.", "✔️ Fix: Weave quotes into sentences (e.g. The author describes the room as 'cavernous')."),
        ("❌ Pitfall 3: Ignoring Punctuation", "Writing long run-on sentences without full stops.", "✔️ Fix: Read your work aloud in your head. Pause where there should be a comma or stop."),
        ("❌ Pitfall 4: Confusing Language & Lit", "Treating stories just as fun tales rather than constructed craft.", "✔️ Fix: Always ask WHY the author chose specific words.")
    ]


    for i, (title, desc, fix) in enumerate(pitfalls):
        top = Inches(1.6 + i * 1.25)
        create_card(slide, Inches(0.8), top, Inches(11.733), Inches(1.1))


        tb = slide.shapes.add_textbox(Inches(1.0), top + Inches(0.1), Inches(11.333), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True


        p0 = tf.paragraphs[0]
        p0.text = title + "  "
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(220, 38, 38) # Red


        run = p0.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = COLOR_DARK_TEXT


        p1 = tf.add_paragraph()
        p1.text = fix
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_ROYAL_BLUE
        p1.space_before = Pt(4)


    add_notes(slide, "TEACHER NOTES (5 mins): Go through common errors seen early in KS3. Establish high standards for sentence structure.")


# SLIDE 20: Real-Life Applications
def build_slide_20():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Real-Life Applications of English Skills", "CONNECTING TO THE REAL WORLD")
    add_footer(slide)


    apps = [
        ("📧 Professional Emails", "Writing clear, polite messages to teachers, bosses, and organizations."),
        ("📰 News & Social Media", "Analyzing articles to spot bias, exaggerated claims, and misdirection."),
        ("🎤 Public Speaking", "Delivering convincing presentations and speeches with vocal projection and confidence."),
        ("📚 Creative Writing", "Crafting compelling stories, scripts, and lyrics that captivate audiences.")
    ]


    for i, (title, desc) in enumerate(apps):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 6.0)
        top = Inches(1.6 + row * 2.6)


        create_card(slide, left, top, Inches(5.733), Inches(2.3))


        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(5.333), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True


        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_ROYAL_BLUE


        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(13)
        p1.font.color.rgb = COLOR_DARK_TEXT
        p1.space_before = Pt(10)


    add_notes(slide, "TEACHER NOTES (4 mins): Remind students that English isn't just an exam subject—it's a life superpower.")


# SLIDE 21: Lesson Summary
def build_slide_21():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Lesson Summary: What Have We Learnt?", "RECAP & SYNTHESIS")
    add_footer(slide)


    summary_items = [
        ("1. Core Distinction", "English Language is the mechanics of communication; English Literature is the study of creative human expression."),
        ("2. Four Pillars", "Mastery of English requires all 4 skills: Reading, Writing, Speaking, and Listening."),
        ("3. Key Terminology", "Author, Audience, Purpose, Genre, and Tone form our foundational vocabulary toolkit."),
        ("4. Text Analysis", "Effective reading involves finding explicit literal information AND making implicit inferences using evidence.")
    ]


    for i, (title, desc) in enumerate(summary_items):
        top = Inches(1.6 + i * 1.25)
        create_card(slide, Inches(0.8), top, Inches(11.733), Inches(1.1))


        tb = slide.shapes.add_textbox(Inches(1.0), top + Inches(0.15), Inches(11.333), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True


        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_ROYAL_BLUE


        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_DARK_TEXT
        p1.space_before = Pt(4)


    add_notes(slide, "TEACHER NOTES (4 mins): Rapid recap. Have individual students summarize one point each.")


# SLIDE 22: Exit Ticket
def build_slide_22():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Exit Ticket: 3-2-1 Reflection", "PLENARY • 5 MINS")
    add_footer(slide)


    prompts = [
        ("3 KEY TERMS", "Write down 3 new vocabulary words you learnt today and their definitions.", COLOR_ROYAL_BLUE),
        ("2 SKILLS", "Name 2 language skills you used during today's lesson activities.", COLOR_ROYAL_BLUE),
        ("1 QUESTION", "Write 1 question you still have about Key Stage 3 English.", COLOR_GOLD)
    ]


    for i, (title, desc, color) in enumerate(prompts):
        left = Inches(0.8 + i * 4.0)
        create_card(slide, left, Inches(1.6), Inches(3.733), Inches(5.0))


        # Header Box
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.6), Inches(3.733), Inches(0.8))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = color
        hdr.line.fill.background()
        tf_h = hdr.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = title
        p_h.font.size = Pt(15)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.alignment = PP_ALIGN.CENTER


        tb = slide.shapes.add_textbox(left + Inches(0.2), Inches(2.6), Inches(3.333), Inches(3.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK_TEXT


    add_notes(slide, "TEACHER NOTES (5 mins): Students complete exit slip on sticky notes or notebook. Collect on way out of classroom.")


# SLIDE 23: Homework
def build_slide_23():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Homework Assignment", "INDEPENDENT STUDY • DUE NEXT LESSON")
    add_footer(slide)


    create_card(slide, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0))
    tb = slide.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.133), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True


    p0 = tf.paragraphs[0]
    p0.text = "📌 Task: Reading & Analysis Log"
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ROYAL_BLUE


    p1 = tf.add_paragraph()
    p1.text = "Read a short story, newspaper article, or book chapter of your choice (minimum 300 words). In your notebook, record:"
    p1.font.size = Pt(14)
    p1.font.color.rgb = COLOR_DARK_TEXT
    p1.space_before = Pt(10)


    tasks = [
        "1. Title and Author of the text.",
        "2. Three NEW or challenging vocabulary words with their dictionary definitions.",
        "3. A 2-sentence summary of the main idea.",
        "4. Copy your FAVOURITE sentence and explain in 1 sentence why you chose it."
    ]
    for t in tasks:
        p = tf.add_paragraph()
        p.text = t
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_ROYAL_BLUE
        p.space_before = Pt(12)


    add_notes(slide, "TEACHER NOTES: Remind students of the due date. Check homework planner logging.")


# SLIDE 24: Preview of Lesson 2
def build_slide_24():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_LIGHT_BG)
    add_header(slide, "Next Time: Preview of Lesson 2", "LOOKING AHEAD")
    add_footer(slide)


    create_card(slide, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0), bg_color=COLOR_ROYAL_BLUE)
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.933), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True


    p0 = tf.paragraphs[0]
    p0.text = "📖 Lesson 2 Topic: Reading Skills — Explicit Info & Simple Inferences"
    p0.font.size = Pt(20)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_GOLD


    p1 = tf.add_paragraph()
    p1.text = "What we will explore next:"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_before = Pt(16)


    previews = [
        "How to skim and scan texts quickly for explicit key facts.",
        "How to use text clues to make accurate, evidence-based inferences.",
        "Analyzing non-fiction news reports vs. fictional adventures.",
        "Mastering P.E.E. paragraphs (Point, Evidence, Explanation)."
    ]
    for pr in previews:
        p = tf.add_paragraph()
        p.text = f"👉 {pr}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_before = Pt(10)


    add_notes(slide, "TEACHER NOTES (2 mins): Build anticipation for next lesson. Encourage students to come prepared.")


# SLIDE 25: Celebration Slide
def build_slide_25():
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, COLOR_ROYAL_BLUE)


    # Accent Gold Border Box
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.5))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_ROYAL_BLUE
    card.line.color.rgb = COLOR_GOLD
    card.line.width = Pt(3)


    tb = slide.shapes.add_textbox(Inches(1.8), Inches(1.3), Inches(9.733), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True


    p0 = tf.paragraphs[0]
    p0.text = "🎉 EXCELLENT WORK TODAY! 🎉"
    p0.font.size = Pt(32)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_GOLD
    p0.alignment = PP_ALIGN.CENTER


    p1 = tf.add_paragraph()
    p1.text = "You have completed Lesson 1 of Key Stage 3 English."
    p1.font.size = Pt(18)
    p1.font.color.rgb = COLOR_WHITE
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(14)


    # Quote Box
    p_q = tf.add_paragraph()
    p_q.text = '"The more that you read, the more things you will know.\nThe more that you learn, the more places you\'ll go."'
    p_q.font.size = Pt(16)
    p_q.font.italic = True
    p_q.font.color.rgb = RGBColor(226, 232, 240)
    p_q.alignment = PP_ALIGN.CENTER
    p_q.space_before = Pt(24)


    p_auth = tf.add_paragraph()
    p_auth.text = "— Dr. Seuss"
    p_auth.font.size = Pt(14)
    p_auth.font.bold = True
    p_auth.font.color.rgb = COLOR_GOLD
    p_auth.alignment = PP_ALIGN.CENTER
    p_auth.space_before = Pt(6)


    p_brand = tf.add_paragraph()
    p_brand.text = FOOTER_TEXT
    p_brand.font.size = Pt(11)
    p_brand.font.color.rgb = RGBColor(148, 163, 184)
    p_brand.alignment = PP_ALIGN.CENTER
    p_brand.space_before = Pt(30)


    add_notes(slide, "TEACHER NOTES: Praise student effort. Wish them a great day and dismiss the class.")


# -----------------------------------------------------------------------------
# MAIN GENERATION PIPELINE
# -----------------------------------------------------------------------------
def build_presentation():
    print("Generating 25 Premium Slides for Bluesky Educational Consults...")
    build_slide_1()
    build_slide_2()
    build_slide_3()
    build_slide_4()
    build_slide_5()
    build_slide_6()
    build_slide_7()
    build_slide_8()
    build_slide_9()
    build_slide_10()
    build_slide_11()
    build_slide_12()
    build_slide_13()
    build_slide_14()
    build_slide_15()
    build_slide_16()
    build_slide_17()
    build_slide_18()
    build_slide_19()
    build_slide_20()
    build_slide_21()
    build_slide_22()
    build_slide_23()
    build_slide_24()
    build_slide_25()


    filename = "Year_7_English_Lesson_1_Bluesky.pptx"
    prs.save(filename)
    print(f"✅ SUCCESS! Presentation saved as '{filename}'.")


if __name__ == "__main__":
    build_presentation()
