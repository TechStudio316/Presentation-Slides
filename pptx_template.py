from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ==========================================
# 1. DESIGN SYSTEM CONSTANTS & STYLES
# ==========================================
NAVY = RGBColor(15, 55, 99)       # #0F3763
TEAL = RGBColor(30, 132, 121)     # #1E8479
LIGHT_GREEN = RGBColor(232, 248, 245) # #E8F8F5
LIGHT_YELLOW = RGBColor(255, 253, 240) # #FFFDF0
YELLOW_ACCENT = RGBColor(255, 191, 0)  # #FFBF00
BG_GRAY = RGBColor(244, 247, 249) # #F4F7F9
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(30, 30, 30)


FONT_NAME = "Arial"




def create_base_presentation():
    """Initializes a 16:9 presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs




def apply_header(slide, brand_text, slide_num, total_slides):
    """Adds the standard header and footer top-bar across slides."""
    # Top Left Brand Tag
    header_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8), Inches(0.4))
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    p.text = brand_text.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_NAME


    # Top Right Slide Tracker
    tracker_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.3), Inches(2.2), Inches(0.4))
    tf2 = tracker_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"Slide {slide_num} of {total_slides}"
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEAL
    p2.font.name = FONT_NAME




def add_pill_badge(slide, text, left, top, width, height, bg_color, text_color):
    """Generates pill-shaped badges (e.g., Week 1, Lesson 1, Duration)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()  # No border
    
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(11)
    p.font.color.rgb = text_color
    p.font.name = FONT_NAME




def add_accent_card(slide, left, top, width, height, accent_color, bg_color=WHITE):
    """Generates a rounded content container with a colored left-accent border bar."""
    # Main Base Box
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = RGBColor(220, 224, 230)
    card.line.width = Pt(1)


    # Left Vertical Accent Line
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()




# ==========================================
# 2. SLIDE GENERATOR FUNCTIONS
# ==========================================


def build_lesson_slide(prs, slide_num, total_slides, title, badges, cards_data):
    """
    Creates a content slide dynamically.
    - badges: list of tuples -> [("WEEK 1", LIGHT_GREEN, TEAL), ...]
    - cards_data: list of dicts -> [{"title": "...", "body": "...", "accent": NAVY}, ...]
    """
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)


    # 1. Header & Brand
    apply_header(slide, "BLUESKY EDUCATIONAL CONSULTS | Learning Beyond Borders", slide_num, total_slides)


    # 2. Render Pill Badges
    badge_x = Inches(0.6)
    for text, bg, color in badges:
        add_pill_badge(slide, text, badge_x, Inches(0.8), Inches(1.5), Inches(0.35), bg, color)
        badge_x += Inches(1.6)


    # 3. Main Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_NAME


    # 4. Render Dynamic Content Cards Grid
    num_cards = len(cards_data)
    if num_cards == 1:
        card_widths = [Inches(12.13)]
        positions = [(Inches(0.6), Inches(2.1))]
    elif num_cards == 2:
        card_widths = [Inches(5.9), Inches(5.9)]
        positions = [(Inches(0.6), Inches(2.1)), (Inches(6.8), Inches(2.1))]
    elif num_cards == 3:
        card_widths = [Inches(3.8), Inches(3.8), Inches(3.8)]
        positions = [(Inches(0.6), Inches(2.1)), (Inches(4.7), Inches(2.1)), (Inches(8.8), Inches(2.1))]
    else:
        # Default 2x2 grid setup
        card_widths = [Inches(5.9), Inches(5.9), Inches(5.9), Inches(5.9)]
        positions = [
            (Inches(0.6), Inches(2.1)), (Inches(6.8), Inches(2.1)),
            (Inches(0.6), Inches(4.7)), (Inches(6.8), Inches(4.7))
        ]


    for idx, card in enumerate(cards_data):
        x, y = positions[idx]
        w = card_widths[idx]
        h = Inches(4.8) if num_cards <= 3 else Inches(2.3)


        # Add visual container
        add_accent_card(slide, x, y, w, h, card.get("accent", NAVY), card.get("bg", WHITE))


        # Add Text Inside Card
        tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), w - Inches(0.5), h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True


        # Card Title
        p_head = tf.paragraphs[0]
        p_head.text = card.get("title", "")
        p_head.font.size = Pt(16)
        p_head.font.bold = True
        p_head.font.color.rgb = card.get("accent", NAVY)
        p_head.font.name = FONT_NAME
        p_head.space_after = Pt(8)


        # Card Content Block
        for line in card.get("bullets", []):
            p_bullet = tf.add_paragraph()
            p_bullet.text = f"•  {line}"
            p_bullet.font.size = Pt(13)
            p_bullet.font.color.rgb = DARK_TEXT
            p_bullet.font.name = FONT_NAME
            p_bullet.space_after = Pt(4)


    return slide




# ==========================================
# 3. EXAMPLE EXECUTION & GENERATION
# ==========================================
if __name__ == "__main__":
    prs = create_base_presentation()


    # Demonstration Slide 1: Learning Objectives (2 Cards Layout)
    build_lesson_slide(
        prs,
        slide_num=4,
        total_slides=25,
        title="Learning Objectives",
        badges=[
            ("WEEK 1", LIGHT_GREEN, TEAL),
            ("LESSON 1", LIGHT_GREEN, TEAL),
            ("90 MINS", LIGHT_YELLOW, YELLOW_ACCENT)
        ],
        cards_data=[
            {
                "title": "Core Grammar Goals",
                "accent": NAVY,
                "bullets": [
                    "Explain clearly what a noun is and state its core function.",
                    "Identify nouns from Key Stage 2 texts.",
                    "Distinguish between common nouns and proper nouns."
                ]
            },
            {
                "title": "Advanced Application",
                "accent": TEAL,
                "bullets": [
                    "Recognise and name groups using collective nouns.",
                    "Identify intangible concepts as abstract nouns.",
                    "Apply all four noun types in creative writing."
                ]
            }
        ]
    )


    # Save generated presentation
    prs.save("Generated_Lesson_Template.pptx")
    print("Presentation saved successfully as 'Generated_Lesson_Template.pptx'.")
