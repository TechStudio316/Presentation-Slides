"""
Year 7 English Lesson 1: Introduction to English Skills
BLUESKY EDUCATIONAL CONSULTS | Learning Beyond Borders
Using the pptx_template.py generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx_template import (
    create_base_presentation, apply_header, add_pill_badge, 
    add_accent_card, build_lesson_slide, 
    NAVY, TEAL, LIGHT_GREEN, LIGHT_YELLOW, YELLOW_ACCENT, 
    BG_GRAY, WHITE, DARK_TEXT, FONT_NAME
)

# Brand Colors
ROYAL_BLUE = RGBColor(10, 78, 163)    # #0A4EA3
GOLD = RGBColor(245, 179, 1)          # #F5B301
LIGHT_GOLD = RGBColor(255, 253, 228)   # Light gold background
SOFT_GREY = RGBColor(244, 247, 249)    # #F4F7F9
MUTED_TEXT = RGBColor(100, 116, 139)   # Slate 500
EMERALD = RGBColor(16, 185, 129)       # Green for answers
RED = RGBColor(220, 38, 38)            # Red for errors

BRAND_TEXT = "BLUESKY EDUCATIONAL CONSULTS | Learning Beyond Borders"


# ==========================================
# CUSTOM HELPER FUNCTIONS
# ==========================================

def set_background(slide, color):
    """Set slide background color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text=BRAND_TEXT):
    """Add footer with divider line."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.0), Inches(12.133), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(220, 224, 230)
    line.line.fill.background()
    
    footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(12.133), Inches(0.3))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED_TEXT
    p.font.name = FONT_NAME


def add_cover_elements(slide, main_title, subtitle, brand=True):
    """Add common cover slide elements."""
    if brand:
        # Brand name at top
        brand_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.4))
        tf = brand_box.text_frame
        p = tf.paragraphs[0]
        p.text = "BLUESKY EDUCATIONAL CONSULTS"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.font.name = FONT_NAME
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = main_title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.8), Inches(12), Inches(0.8))
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = GOLD
        p.font.name = FONT_NAME
    
    add_footer(slide)


def add_question_card(slide, left, top, width, height, title, question, answer=None):
    """Add a question card with optional answer."""
    add_accent_card(slide, left, top, width, height, NAVY, WHITE)
    
    tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2), width - Inches(0.5), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(14)
    p_title.font.bold = True
    p_title.font.color.rgb = NAVY
    p_title.font.name = FONT_NAME
    
    p_q = tf.add_paragraph()
    p_q.text = question
    p_q.font.size = Pt(12)
    p_q.font.color.rgb = DARK_TEXT
    p_q.space_before = Pt(8)
    
    if answer:
        p_a = tf.add_paragraph()
        p_a.text = f"✓ {answer}"
        p_a.font.size = Pt(12)
        p_a.font.bold = True
        p_a.font.color.rgb = EMERALD
        p_a.space_before = Pt(8)


# ==========================================
# SLIDE BUILDERS
# ==========================================

def build_slide_1_cover(prs):
    """Slide 1: Premium Cover Slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, ROYAL_BLUE)
    
    # Gold accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(0.15), Inches(5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11), Inches(0.5))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = "BLUESKY EDUCATIONAL CONSULTS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = FONT_NAME
    
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.5), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Introduction to English Skills"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11), Inches(1))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Understanding Language, Literature & Effective Communication"
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(226, 232, 240)
    p.font.name = FONT_NAME
    
    # Details
    details_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11), Inches(1))
    tf = details_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Year 7 English • Key Stage 3 • Week 1, Lesson 1"
    p.font.size = Pt(14)
    p.font.color.rgb = GOLD
    p.font.name = FONT_NAME
    
    p2 = tf.add_paragraph()
    p2.text = "Duration: 90 Minutes"
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(226, 232, 240)
    p2.space_before = Pt(8)
    
    add_footer(slide)
    return slide


def build_slide_2_welcome(prs, total):
    """Slide 2: Welcome Back."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 2, total)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Welcome to Key Stage 3 English!"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    # Main welcome card
    add_accent_card(slide, Inches(0.6), Inches(1.9), Inches(7.5), Inches(4.8), ROYAL_BLUE, WHITE)
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(6.9), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "A New Chapter Begins"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    
    bullets = [
        "Welcome to your Secondary English journey!",
        "Explore powerful stories, debate big ideas, and master expression.",
        "Our classroom is a safe space to test ideas and grow.",
        "Read critically, write creatively, speak persuasively."
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(12)
    
    # Values card
    add_accent_card(slide, Inches(8.4), Inches(1.9), Inches(4.3), Inches(4.8), GOLD, ROYAL_BLUE)
    tb2 = slide.shapes.add_textbox(Inches(8.7), Inches(2.1), Inches(3.8), Inches(4.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    p2 = tf2.paragraphs[0]
    p2.text = "Classroom Values"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = GOLD
    
    values = ["📖 Active Listening", "🙋 Active Participation", "🤝 Mutual Respect", "💡 Creative Curiosity", "✏️ High Endeavour"]
    for v in values:
        p = tf2.add_paragraph()
        p.text = v
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_before = Pt(14)
    
    add_footer(slide)
    return slide


def build_slide_3_icebreaker(prs, total):
    """Slide 3: Ice Breaker."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 3, total)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ice Breaker: What Makes a Good Communicator?"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    # Question card
    add_accent_card(slide, Inches(0.6), Inches(1.8), Inches(12.133), Inches(1.5), GOLD, ROYAL_BLUE)
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.95), Inches(11.6), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💬 Think of someone you love listening to (a YouTuber, teacher, author, or friend)."
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "What makes them so effective when they communicate?"
    p2.font.size = Pt(14)
    p2.font.color.rgb = GOLD
    p2.space_before = Pt(6)
    
    # Discussion prompts
    prompts = [
        ("1. Tone & Voice", "Is it how they sound, the words they choose, or their body language?"),
        ("2. Clarity", "Do they make complex ideas easy to understand?"),
        ("3. Engagement", "How do they keep your attention from start to finish?")
    ]
    
    for i, (title, desc) in enumerate(prompts):
        left = Inches(0.6 + i * 4.1)
        add_accent_card(slide, left, Inches(3.5), Inches(3.8), Inches(3.2), ROYAL_BLUE, WHITE)
        
        tb_p = slide.shapes.add_textbox(left + Inches(0.25), Inches(3.7), Inches(3.3), Inches(2.8))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True
        
        p_t = tf_p.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = ROYAL_BLUE
        
        p_d = tf_p.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = DARK_TEXT
        p_d.space_before = Pt(10)
    
    add_footer(slide)
    return slide


def build_slide_4_objectives(prs, total):
    """Slide 4: Learning Objectives."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 4, total)
    
    # Badges
    badges = [("WEEK 1", LIGHT_GREEN, TEAL), ("LESSON 1", LIGHT_GREEN, TEAL), ("90 MINS", LIGHT_YELLOW, RGBColor(200, 140, 0))]
    badge_x = Inches(0.6)
    for text, bg, color in badges:
        add_pill_badge(slide, text, badge_x, Inches(0.85), Inches(1.4), Inches(0.32), bg, color)
        badge_x += Inches(1.5)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Learning Objectives"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    objectives = [
        ("01", "Distinguish Core Domains", "Explain what English Language and English Literature are."),
        ("02", "Identify the Four Pillars", "Recognise the four key skills: Reading, Writing, Speaking, Listening."),
        ("03", "Understand Communication", "Recognise why effective communication is vital."),
        ("04", "Master Key Vocabulary", "Use subject-specific terminology confidently."),
        ("05", "Analyse Text for Meaning", "Extract literal and inferential meaning using evidence.")
    ]
    
    for i, (num, title, desc) in enumerate(objectives):
        top = Inches(1.9 + i * 0.95)
        
        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), top, Inches(0.7), Inches(0.7))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ROYAL_BLUE
        badge.line.fill.background()
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = num
        p_b.font.size = Pt(14)
        p_b.font.bold = True
        p_b.font.color.rgb = GOLD
        p_b.alignment = PP_ALIGN.CENTER
        
        # Content card
        add_accent_card(slide, Inches(1.5), top, Inches(11.233), Inches(0.75), ROYAL_BLUE, WHITE)
        tb = slide.shapes.add_textbox(Inches(1.7), top + Inches(0.1), Inches(10.8), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_t = tf.paragraphs[0]
        p_t.text = title + " — "
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = ROYAL_BLUE
        
        run = p_t.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = DARK_TEXT
    
    add_footer(slide)
    return slide


def build_slide_5_success(prs, total):
    """Slide 5: Success Criteria."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 5, total)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Success Criteria"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    criteria = [
        ("Bronze Level", "I can state the difference between English Language and English Literature, and list the four skills.", ROYAL_BLUE),
        ("Silver Level", "I can define key vocabulary and identify them in a short text.", TEAL),
        ("Gold Level", "I can explain the author's purpose using explicit evidence and make simple inferences.", GOLD)
    ]
    
    for i, (level, desc, color) in enumerate(criteria):
        top = Inches(1.8 + i * 1.6)
        add_accent_card(slide, Inches(0.6), top, Inches(12.133), Inches(1.3), color, WHITE)
        
        tb = slide.shapes.add_textbox(Inches(0.9), top + Inches(0.15), Inches(11.6), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = f"✓ {level}"
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = color
        
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(14)
        p1.font.color.rgb = DARK_TEXT
        p1.space_before = Pt(6)
    
    add_footer(slide)
    return slide


def build_slide_6_warmup(prs, total):
    """Slide 6: Warm-up Activity."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 6, total)
    
    # Badges
    badges = [("STARTER", LIGHT_GREEN, TEAL), ("5 MINS", LIGHT_YELLOW, RGBColor(200, 140, 0))]
    badge_x = Inches(0.6)
    for text, bg, color in badges:
        add_pill_badge(slide, text, badge_x, Inches(0.85), Inches(1.3), Inches(0.32), bg, color)
        badge_x += Inches(1.4)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Warm-up: Word Association Challenge"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    # Instruction banner
    add_accent_card(slide, Inches(0.6), Inches(1.9), Inches(12.133), Inches(0.8), GOLD, ROYAL_BLUE)
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.6), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "⚡ Challenge: You have 2 minutes! Write down words you associate with each pillar:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    pillars = [
        ("📖 READING", ["Books", "Comprehension", "Library", "Context", "Analyze", "Decode"]),
        ("✍️ WRITING", ["Grammar", "Essays", "Creativity", "Vocabulary", "Draft", "Punctuation"]),
        ("🗣️ SPEAKING", ["Debate", "Tone", "Pitch", "Expression", "Dialogue", "Presentation"]),
        ("👂 LISTENING", ["Attention", "Empathy", "Focus", "Active", "Notes", "Understanding"])
    ]
    
    for i, (title, examples) in enumerate(pillars):
        left = Inches(0.6 + i * 3.1)
        add_accent_card(slide, left, Inches(2.9), Inches(2.9), Inches(3.8), ROYAL_BLUE, WHITE)
        
        tb_c = slide.shapes.add_textbox(left + Inches(0.2), Inches(3.05), Inches(2.5), Inches(3.5))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        
        p_t = tf_c.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = ROYAL_BLUE
        
        p_sub = tf_c.add_paragraph()
        p_sub.text = "Examples:"
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = GOLD
        p_sub.space_before = Pt(8)
        
        for ex in examples:
            p_ex = tf_c.add_paragraph()
            p_ex.text = f"• {ex}"
            p_ex.font.size = Pt(11)
            p_ex.font.color.rgb = DARK_TEXT
            p_ex.space_before = Pt(3)
    
    add_footer(slide)
    return slide


def build_slide_7_what_is_english(prs, total):
    """Slide 7: What is English?"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 7, total)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "What is English? Language vs. Literature"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    # Language card
    add_accent_card(slide, Inches(0.6), Inches(1.8), Inches(6.0), Inches(4.8), ROYAL_BLUE, WHITE)
    tb1 = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(5.4), Inches(4.4))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "🔤 English Language"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = ROYAL_BLUE
    
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "The Tool of Communication"
    p1_sub.font.size = Pt(12)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = GOLD
    p1_sub.space_before = Pt(4)
    
    lang_points = [
        "Focuses on HOW words work together.",
        "Includes grammar, punctuation, vocabulary.",
        "Changes according to audience & purpose.",
        "Teaches non-fiction: articles, speeches, reports."
    ]
    for pt in lang_points:
        p = tf1.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(10)
    
    # Literature card
    add_accent_card(slide, Inches(6.9), Inches(1.8), Inches(6.0), Inches(4.8), GOLD, WHITE)
    tb2 = slide.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(5.4), Inches(4.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    
    p2 = tf2.paragraphs[0]
    p2.text = "📚 English Literature"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = ROYAL_BLUE
    
    p2_sub = tf2.add_paragraph()
    p2_sub.text = "The Study of Human Experience"
    p2_sub.font.size = Pt(12)
    p2_sub.font.bold = True
    p2_sub.font.color.rgb = GOLD
    p2_sub.space_before = Pt(4)
    
    lit_points = [
        "Focuses on WHAT writers express and WHY.",
        "Includes novels, plays, poetry, historical texts.",
        "Explores themes: love, power, conflict, identity.",
        "Teaches analytical and critical thinking."
    ]
    for pt in lit_points:
        p = tf2.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(10)
    
    add_footer(slide)
    return slide


def build_slide_8_four_skills(prs, total):
    """Slide 8: The Four Language Skills."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 8, total)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "The Four Language Skills"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    skills = [
        ("📖 READING", "Receptive", "Decoding symbols, interpreting meaning, analyzing techniques.", ROYAL_BLUE),
        ("✍️ WRITING", "Productive", "Structuring thoughts, selecting vocabulary, crafting narratives.", ROYAL_BLUE),
        ("🗣️ SPEAKING", "Productive", "Articulating ideas, adapting tone, presenting arguments.", GOLD),
        ("👂 LISTENING", "Receptive", "Processing spoken words, active listening, understanding intent.", GOLD)
    ]
    
    for i, (title, type_label, desc, color) in enumerate(skills):
        row = i // 2
        col = i % 2
        left = Inches(0.6 + col * 6.3)
        top = Inches(1.8 + row * 2.5)
        
        add_accent_card(slide, left, top, Inches(6.0), Inches(2.2), color, WHITE)
        
        tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), Inches(5.5), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(18)
        p0.font.bold = True
        p0.font.color.rgb = color
        
        p_sub = tf.add_paragraph()
        p_sub.text = f"[{type_label} Skill]"
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = MUTED_TEXT
        
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(13)
        p1.font.color.rgb = DARK_TEXT
        p1.space_before = Pt(8)
    
    add_footer(slide)
    return slide


def build_slide_9_why_english(prs, total):
    """Slide 9: Why English Matters."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 9, total)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Why English Matters Beyond the Classroom"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    areas = [
        ("🎓 Education", "Mastering English unlocks all subjects—from science reports to history documents."),
        ("💼 Careers", "Employers rank clear communication as the #1 essential workplace skill."),
        ("🌐 Media Literacy", "Helps navigate social media, spot fake news, evaluate bias."),
        ("🗣️ Everyday Life", "Expressing feelings, resolving conflicts, forming relationships.")
    ]
    
    for i, (title, desc) in enumerate(areas):
        left = Inches(0.6 + i * 3.1)
        add_accent_card(slide, left, Inches(1.8), Inches(2.9), Inches(4.8), ROYAL_BLUE, WHITE)
        
        tb = slide.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(2.5), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = ROYAL_BLUE
        
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(12)
        p1.font.color.rgb = DARK_TEXT
        p1.space_before = Pt(12)
    
    add_footer(slide)
    return slide


def build_slide_10_vocabulary(prs, total):
    """Slide 10: Building Vocabulary."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 10, total)
    
    # Badges
    badges = [("VOCABULARY", LIGHT_GREEN, TEAL)]
    add_pill_badge(slide, badges[0][0], Inches(0.6), Inches(0.85), Inches(1.8), Inches(0.32), badges[0][1], badges[0][2])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Building Vocabulary: Core English Toolkit"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    terms = [
        ("Author", "The writer of a text, story, or document."),
        ("Audience", "The intended reader, listener, or viewer."),
        ("Purpose", "The reason a text was created (inform, persuade, entertain)."),
        ("Genre", "The category or style of literature/media."),
        ("Tone", "The mood or attitude expressed by the writer.")
    ]
    
    for i, (term, defn) in enumerate(terms):
        top = Inches(1.9 + i * 0.95)
        
        # Term header
        add_accent_card(slide, Inches(0.6), top, Inches(2.3), Inches(0.75), ROYAL_BLUE, ROYAL_BLUE)
        tb_t = slide.shapes.add_textbox(Inches(0.75), top + Inches(0.15), Inches(2.0), Inches(0.45))
        p_t = tb_t.text_frame.paragraphs[0]
        p_t.text = term
        p_t.font.size = Pt(14)
        p_t.font.bold = True
        p_t.font.color.rgb = GOLD
        
        # Definition
        add_accent_card(slide, Inches(3.1), top, Inches(9.633), Inches(0.75), RGBColor(220, 224, 230), WHITE)
        tb_d = slide.shapes.add_textbox(Inches(3.3), top + Inches(0.15), Inches(9.2), Inches(0.45))
        p_d = tb_d.text_frame.paragraphs[0]
        p_d.text = defn
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = DARK_TEXT
    
    add_footer(slide)
    return slide


def build_slide_11_reading(prs, total):
    """Slide 11: Reading Activity."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 11, total)
    
    # Badges
    badges = [("READING", LIGHT_GREEN, TEAL), ("10 MINS", LIGHT_YELLOW, RGBColor(200, 140, 0))]
    badge_x = Inches(0.6)
    for text, bg, color in badges:
        add_pill_badge(slide, text, badge_x, Inches(0.85), Inches(1.3), Inches(0.32), bg, color)
        badge_x += Inches(1.4)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Reading Activity: The Discovery"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    # Passage card
    add_accent_card(slide, Inches(0.6), Inches(1.9), Inches(12.133), Inches(4.0), RGBColor(220, 224, 230), WHITE)
    
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.6), Inches(3.6))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_head = tf.paragraphs[0]
    p_head.text = "Excerpt from 'The Forgotten Library' by Maya Lin"
    p_head.font.size = Pt(12)
    p_head.font.italic = True
    p_head.font.color.rgb = MUTED_TEXT
    
    passage = (
        "The heavy wooden door groaned as Arthur pushed it open, revealing a cavernous room "
        "bathed in golden dust motes. Bookshelves reached up toward vaulted ceilings like ancient oak trees, "
        "their branches heavy with leather-bound volumes. A gentle smell of old paper and dried lavender "
        "hung in the still air. In the centre stood a polished mahogany desk, and upon it lay "
        "a single brass key resting on an unopened letter. Arthur took a slow step forward, his heart "
        "hammering against his ribs like a trapped bird."
    )
    
    p_body = tf.add_paragraph()
    p_body.text = passage
    p_body.font.size = Pt(14)
    p_body.font.color.rgb = DARK_TEXT
    p_body.space_before = Pt(12)
    
    # Task banner
    add_accent_card(slide, Inches(0.6), Inches(6.1), Inches(12.133), Inches(0.65), GOLD, ROYAL_BLUE)
    tb_c = slide.shapes.add_textbox(Inches(0.9), Inches(6.2), Inches(11.6), Inches(0.45))
    p_c = tb_c.text_frame.paragraphs[0]
    p_c.text = "🔍 Task: Read twice. Notice how the author creates atmosphere using sensory details."
    p_c.font.size = Pt(13)
    p_c.font.bold = True
    p_c.font.color.rgb = WHITE
    
    add_footer(slide)
    return slide


def build_slide_12_modelling(prs, total):
    """Slide 12: Teacher Modelling."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 12, total)
    
    # Badges
    badges = [("I DO", LIGHT_GREEN, TEAL), ("ANNOTATING", LIGHT_YELLOW, RGBColor(200, 140, 0))]
    badge_x = Inches(0.6)
    for text, bg, color in badges:
        add_pill_badge(slide, text, badge_x, Inches(0.85), Inches(1.4), Inches(0.32), bg, color)
        badge_x += Inches(1.5)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Teacher Modelling: Think Aloud Strategy"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    # Text focus card
    add_accent_card(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(4.8), RGBColor(220, 224, 230), WHITE)
    tb_l = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(5.2), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    
    p0 = tf_l.paragraphs[0]
    p0.text = "Text Focus:"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = GOLD
    
    p1 = tf_l.add_paragraph()
    p1.text = '"...his heart hammering against his ribs like a trapped bird."'
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = ROYAL_BLUE
    p1.space_before = Pt(10)
    
    # Think aloud card
    add_accent_card(slide, Inches(6.7), Inches(1.9), Inches(5.993), Inches(4.8), ROYAL_BLUE, ROYAL_BLUE)
    tb_r = slide.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.4), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    
    p2 = tf_r.paragraphs[0]
    p2.text = "🧠 Teacher's 'Think Aloud':"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = GOLD
    
    thoughts = [
        "1. Identify Technique: 'like a trapped bird' is a Simile.",
        "2. Analyze Imagery: A trapped bird flutters frantically.",
        "3. Connect to Meaning: Arthur feels overwhelmed by mystery.",
        "4. Evaluate Atmosphere: The writer creates tension."
    ]
    for th in thoughts:
        p = tf_r.add_paragraph()
        p.text = th
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)
    
    add_footer(slide)
    return slide


def build_slide_13_guided(prs, total):
    """Slide 13: Guided Reading."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 13, total)
    
    # Badges
    badges = [("WE DO", LIGHT_GREEN, TEAL), ("GUIDED", LIGHT_YELLOW, RGBColor(200, 140, 0))]
    badge_x = Inches(0.6)
    for text, bg, color in badges:
        add_pill_badge(slide, text, badge_x, Inches(0.85), Inches(1.3), Inches(0.32), bg, color)
        badge_x += Inches(1.4)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Guided Reading: Let's Unpack Together"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    questions = [
        ("Q1: Sensory Language", "Which senses does the author appeal to in line 3?", "Smell (old paper, lavender) and Sound (door groaned)."),
        ("Q2: Figurative Language", "What metaphor describes the bookshelves?", "Compared to 'ancient oak trees' with branches heavy with books."),
        ("Q3: Mystery & Suspense", "Which object signals a story is about to begin?", "The brass key resting on the unopened letter.")
    ]
    
    for i, (q_title, q_text, q_ans) in enumerate(questions):
        top = Inches(1.9 + i * 1.65)
        add_accent_card(slide, Inches(0.6), top, Inches(12.133), Inches(1.45), RGBColor(220, 224, 230), WHITE)
        
        tb = slide.shapes.add_textbox(Inches(0.9), top + Inches(0.15), Inches(11.6), Inches(1.15))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = q_title
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = ROYAL_BLUE
        
        p1 = tf.add_paragraph()
        p1.text = q_text
        p1.font.size = Pt(12)
        p1.font.color.rgb = DARK_TEXT
        
        p2 = tf.add_paragraph()
        p2.text = f"💡 Answer: {q_ans}"
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = EMERALD
        p2.space_before = Pt(4)
    
    add_footer(slide)
    return slide


def build_slide_14_comprehension(prs, total):
    """Slide 14: Comprehension Questions."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 14, total)
    
    # Badges
    badges = [("YOU DO", LIGHT_GREEN, TEAL)]
    add_pill_badge(slide, badges[0][0], Inches(0.6), Inches(0.85), Inches(1.2), Inches(0.32), badges[0][1], badges[0][2])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Comprehension: Literal vs. Inferential"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    # Literal questions
    add_accent_card(slide, Inches(0.6), Inches(1.9), Inches(6.0), Inches(4.8), RGBColor(220, 224, 230), WHITE)
    tb_l = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(5.4), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    
    p_l_head = tf_l.paragraphs[0]
    p_l_head.text = "📍 Literal Questions (In the text)"
    p_l_head.font.size = Pt(14)
    p_l_head.font.bold = True
    p_l_head.font.color.rgb = ROYAL_BLUE
    
    l_qs = [
        "1. What kind of door did Arthur push open?\n   → A heavy wooden door.",
        "2. What object was on the mahogany desk?\n   → A brass key and an unopened letter."
    ]
    for q in l_qs:
        p = tf_l.add_paragraph()
        p.text = q
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(14)
    
    # Inferential questions
    add_accent_card(slide, Inches(6.9), Inches(1.9), Inches(6.0), Inches(4.8), GOLD, WHITE)
    tb_r = slide.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.4), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    
    p_r_head = tf_r.paragraphs[0]
    p_r_head.text = "🔍 Inferential (Read between lines)"
    p_r_head.font.size = Pt(14)
    p_r_head.font.bold = True
    p_r_head.font.color.rgb = ROYAL_BLUE
    
    r_qs = [
        "3. Has anyone been in this library recently?\n   → No. Dust motes and still air suggest untouched.",
        "4. How does Arthur feel about being there?\n   → Both awed and terrified (heart hammering)."
    ]
    for q in r_qs:
        p = tf_r.add_paragraph()
        p.text = q
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(14)
    
    add_footer(slide)
    return slide


def build_slide_15_vocab_activity(prs, total):
    """Slide 15: Interactive Vocabulary Activity."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 15, total)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Vocabulary Matching Challenge"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    # Terms card
    add_accent_card(slide, Inches(0.6), Inches(1.7), Inches(4.2), Inches(4.9), ROYAL_BLUE, ROYAL_BLUE)
    tb_l = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(3.6), Inches(4.5))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.paragraphs[0].text = "TERMS"
    tf_l.paragraphs[0].font.size = Pt(14)
    tf_l.paragraphs[0].font.bold = True
    tf_l.paragraphs[0].font.color.rgb = GOLD
    
    terms_list = ["1. Cavernous", "2. Vaulted", "3. Hammering", "4. Stillness"]
    for t in terms_list:
        p = tf_l.add_paragraph()
        p.text = t
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_before = Pt(18)
    
    # Definitions card
    add_accent_card(slide, Inches(5.1), Inches(1.7), Inches(7.633), Inches(4.9), RGBColor(220, 224, 230), WHITE)
    tb_r = slide.shapes.add_textbox(Inches(5.4), Inches(1.9), Inches(7.0), Inches(4.5))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.paragraphs[0].text = "DEFINITIONS"
    tf_r.paragraphs[0].font.size = Pt(14)
    tf_r.paragraphs[0].font.bold = True
    tf_r.paragraphs[0].font.color.rgb = ROYAL_BLUE
    
    defs_list = [
        "A. Beating rapidly and loudly  ➔ [#3]",
        "B. Vast, huge, dark like a cave  ➔ [#1]",
        "C. Quiet, calm, undisturbed  ➔ [#4]",
        "D. High, arched ceilings  ➔ [#2]"
    ]
    for d in defs_list:
        p = tf_r.add_paragraph()
        p.text = d
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(16)
    
    add_footer(slide)
    return slide


def build_slide_16_think_pair_share(prs, total):
    """Slide 16: Think-Pair-Share."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 16, total)
    
    # Badges
    badges = [("COLLABORATIVE", LIGHT_GREEN, TEAL), ("6 MINS", LIGHT_YELLOW, RGBColor(200, 140, 0))]
    badge_x = Inches(0.6)
    for text, bg, color in badges:
        add_pill_badge(slide, text, badge_x, Inches(0.85), Inches(1.6), Inches(0.32), bg, color)
        badge_x += Inches(1.7)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Think–Pair–Share: Author's Purpose"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    steps = [
        ("THINK (1 Min)", "Consider: Why did Maya Lin write this passage? What emotion did she want you to feel?", ROYAL_BLUE),
        ("PAIR (2 Mins)", "Compare ideas with your partner. Did you pick the same words? Agree on the main purpose.", TEAL),
        ("SHARE (3 Mins)", "Be ready to share your partner's best idea with evidence from the text.", GOLD)
    ]
    
    for i, (title, desc, color) in enumerate(steps):
        left = Inches(0.6 + i * 4.15)
        add_accent_card(slide, left, Inches(1.9), Inches(3.9), Inches(4.7), color, WHITE)
        
        tb = slide.shapes.add_textbox(left + Inches(0.25), Inches(2.1), Inches(3.4), Inches(4.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = color
        
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(12)
        p1.font.color.rgb = DARK_TEXT
        p1.space_before = Pt(12)
    
    add_footer(slide)
    return slide


def build_slide_17_independent(prs, total):
    """Slide 17: Independent Practice."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 17, total)
    
    # Badges
    badges = [("YOU DO", LIGHT_GREEN, TEAL), ("INDIVIDUAL", LIGHT_YELLOW, RGBColor(200, 140, 0)), ("10 MINS", RGBColor(255, 230, 230), RED)]
    badge_x = Inches(0.6)
    for text, bg, color in badges:
        add_pill_badge(slide, text, badge_x, Inches(0.85), Inches(1.5), Inches(0.32), bg, color)
        badge_x += Inches(1.6)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Independent Practice: Written Analysis"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    # Task card
    add_accent_card(slide, Inches(0.6), Inches(1.9), Inches(12.133), Inches(4.8), RGBColor(220, 224, 230), WHITE)
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "📝 Task: Answer in full sentences in your notebook:"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = ROYAL_BLUE
    
    p_prompt = tf.add_paragraph()
    p_prompt.text = '"How does Maya Lin create a sense of mystery and wonder?"'
    p_prompt.font.size = Pt(16)
    p_prompt.font.bold = True
    p_prompt.font.color.rgb = GOLD
    p_prompt.space_before = Pt(10)
    
    p_guide = tf.add_paragraph()
    p_guide.text = "Requirements:"
    p_guide.font.size = Pt(12)
    p_guide.font.bold = True
    p_guide.font.color.rgb = DARK_TEXT
    p_guide.space_before = Pt(14)
    
    reqs = [
        "Write 3-4 full sentences.",
        "Include ONE direct quotation from the text.",
        "Use key vocabulary (Tone, Imagery, Audience, Atmosphere).",
        "Check capitals, full stops, and spelling."
    ]
    for r in reqs:
        p = tf.add_paragraph()
        p.text = f"✓ {r}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(6)
    
    add_footer(slide)
    return slide


def build_slide_18_creative(prs, total):
    """Slide 18: Creative Challenge."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 18, total)
    
    # Badges
    badges = [("CREATIVE", LIGHT_GREEN, TEAL), ("8 MINS", LIGHT_YELLOW, RGBColor(200, 140, 0))]
    badge_x = Inches(0.6)
    for text, bg, color in badges:
        add_pill_badge(slide, text, badge_x, Inches(0.85), Inches(1.3), Inches(0.32), bg, color)
        badge_x += Inches(1.4)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Creative Writing: Self-Introduction"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    # Task card
    add_accent_card(slide, Inches(0.6), Inches(1.9), Inches(12.133), Inches(4.8), GOLD, WHITE)
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "🎨 Task: Introduce Yourself Using Descriptive Language"
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = ROYAL_BLUE
    
    p1 = tf.add_paragraph()
    p1.text = "Write a 50-word paragraph introducing yourself to the class:"
    p1.font.size = Pt(14)
    p1.font.color.rgb = DARK_TEXT
    p1.space_before = Pt(10)
    
    rules = [
        "Do NOT just list facts (e.g. 'My name is Alex and I like football').",
        "Use vivid verbs and sensory descriptions.",
        "Include ONE simile or metaphor for your personality.",
        "Set an enthusiastic or mysterious TONE."
    ]
    for r in rules:
        p = tf.add_paragraph()
        p.text = f"⭐ {r}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(10)
    
    add_footer(slide)
    return slide


def build_slide_19_mistakes(prs, total):
    """Slide 19: Common Mistakes."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 19, total)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Common Pitfalls in Year 7 English"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    pitfalls = [
        ("❌ One-Word Answers", "Giving 'Yes' or 'It was good'.", "✓ Always answer in full sentences with reasons."),
        ("❌ Floating Quotations", "Dropping quotes without context.", "✓ Weave quotes into sentences."),
        ("❌ Ignoring Punctuation", "Run-on sentences without stops.", "✓ Read aloud and pause where needed."),
        ("❌ Confusing Language & Lit", "Treating stories as just fun tales.", "✓ Always ask WHY the author chose words.")
    ]
    
    for i, (title, desc, fix) in enumerate(pitfalls):
        top = Inches(1.7 + i * 1.25)
        add_accent_card(slide, Inches(0.6), top, Inches(12.133), Inches(1.05), RGBColor(220, 224, 230), WHITE)
        
        tb = slide.shapes.add_textbox(Inches(0.9), top + Inches(0.1), Inches(11.6), Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = title + "  "
        p0.font.size = Pt(12)
        p0.font.bold = True
        p0.font.color.rgb = RED
        
        run = p0.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = DARK_TEXT
        
        p1 = tf.add_paragraph()
        p1.text = fix
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = EMERALD
        p1.space_before = Pt(4)
    
    add_footer(slide)
    return slide


def build_slide_20_real_life(prs, total):
    """Slide 20: Real-Life Applications."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 20, total)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Real-Life Applications of English Skills"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    apps = [
        ("📧 Professional Emails", "Writing clear, polite messages to teachers, bosses, organisations."),
        ("📰 News & Social Media", "Analyzing articles to spot bias and misdirection."),
        ("🎤 Public Speaking", "Delivering convincing presentations with confidence."),
        ("📚 Creative Writing", "Crafting compelling stories and scripts.")
    ]
    
    for i, (title, desc) in enumerate(apps):
        row = i // 2
        col = i % 2
        left = Inches(0.6 + col * 6.3)
        top = Inches(1.7 + row * 2.5)
        
        add_accent_card(slide, left, top, Inches(6.0), Inches(2.2), ROYAL_BLUE, WHITE)
        
        tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), Inches(5.5), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = ROYAL_BLUE
        
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(13)
        p1.font.color.rgb = DARK_TEXT
        p1.space_before = Pt(10)
    
    add_footer(slide)
    return slide


def build_slide_21_summary(prs, total):
    """Slide 21: Lesson Summary."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 21, total)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Lesson Summary: What Have We Learnt?"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    summary = [
        ("1. Core Distinction", "English Language is mechanics; English Literature is creative expression."),
        ("2. Four Pillars", "Mastery requires all 4 skills: Reading, Writing, Speaking, Listening."),
        ("3. Key Terminology", "Author, Audience, Purpose, Genre, and Tone form our toolkit."),
        ("4. Text Analysis", "Find literal info AND make inferences using evidence.")
    ]
    
    for i, (title, desc) in enumerate(summary):
        top = Inches(1.7 + i * 1.25)
        add_accent_card(slide, Inches(0.6), top, Inches(12.133), Inches(1.05), RGBColor(220, 224, 230), WHITE)
        
        tb = slide.shapes.add_textbox(Inches(0.9), top + Inches(0.15), Inches(11.6), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(14)
        p0.font.bold = True
        p0.font.color.rgb = ROYAL_BLUE
        
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(12)
        p1.font.color.rgb = DARK_TEXT
        p1.space_before = Pt(4)
    
    add_footer(slide)
    return slide


def build_slide_22_exit(prs, total):
    """Slide 22: Exit Ticket."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 22, total)
    
    # Badges
    badges = [("PLENARY", LIGHT_GREEN, TEAL), ("5 MINS", LIGHT_YELLOW, RGBColor(200, 140, 0))]
    badge_x = Inches(0.6)
    for text, bg, color in badges:
        add_pill_badge(slide, text, badge_x, Inches(0.85), Inches(1.3), Inches(0.32), bg, color)
        badge_x += Inches(1.4)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Exit Ticket: 3-2-1 Reflection"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    prompts = [
        ("3 KEY TERMS", "Write 3 new vocabulary words and their definitions.", ROYAL_BLUE),
        ("2 SKILLS", "Name 2 language skills you used today.", TEAL),
        ("1 QUESTION", "Write 1 question you still have about KS3 English.", GOLD)
    ]
    
    for i, (title, desc, color) in enumerate(prompts):
        left = Inches(0.6 + i * 4.15)
        add_accent_card(slide, left, Inches(1.9), Inches(3.9), Inches(4.7), color, WHITE)
        
        tb = slide.shapes.add_textbox(left + Inches(0.25), Inches(2.1), Inches(3.4), Inches(4.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(15)
        p0.font.bold = True
        p0.font.color.rgb = color
        
        p1 = tf.add_paragraph()
        p1.text = desc
        p1.font.size = Pt(12)
        p1.font.color.rgb = DARK_TEXT
        p1.space_before = Pt(12)
    
    add_footer(slide)
    return slide


def build_slide_23_homework(prs, total):
    """Slide 23: Homework."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 23, total)
    
    # Badges
    badges = [("HOMEWORK", RGBColor(255, 230, 230), RED)]
    add_pill_badge(slide, badges[0][0], Inches(0.6), Inches(0.85), Inches(1.5), Inches(0.32), badges[0][1], badges[0][2])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Homework Assignment"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    # Task card
    add_accent_card(slide, Inches(0.6), Inches(1.9), Inches(12.133), Inches(4.8), RGBColor(220, 224, 230), WHITE)
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "📌 Task: Reading & Analysis Log"
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = ROYAL_BLUE
    
    p1 = tf.add_paragraph()
    p1.text = "Read a short story or article (min 300 words). Record:"
    p1.font.size = Pt(14)
    p1.font.color.rgb = DARK_TEXT
    p1.space_before = Pt(10)
    
    tasks = [
        "1. Title and Author of the text.",
        "2. Three NEW vocabulary words with definitions.",
        "3. A 2-sentence summary of the main idea.",
        "4. Your FAVOURITE sentence and why you chose it."
    ]
    for t in tasks:
        p = tf.add_paragraph()
        p.text = t
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = ROYAL_BLUE
        p.space_before = Pt(12)
    
    add_footer(slide)
    return slide


def build_slide_24_preview(prs, total):
    """Slide 24: Preview of Lesson 2."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, SOFT_GREY)
    apply_header(slide, BRAND_TEXT, 24, total)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Next Time: Preview of Lesson 2"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.font.name = FONT_NAME
    
    # Preview card
    add_accent_card(slide, Inches(0.6), Inches(1.7), Inches(12.133), Inches(5.0), ROYAL_BLUE, ROYAL_BLUE)
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.6), Inches(4.6))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "📖 Lesson 2: Reading Skills"
    p0.font.size = Pt(20)
    p0.font.bold = True
    p0.font.color.rgb = GOLD
    
    p1 = tf.add_paragraph()
    p1.text = "Finding Explicit Information and Making Simple Inferences"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_before = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = "What we will explore:"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = GOLD
    p2.space_before = Pt(16)
    
    previews = [
        "How to skim and scan texts quickly for key facts.",
        "How to use text clues to make evidence-based inferences.",
        "Analyzing non-fiction vs. fictional texts.",
        "Mastering P.E.E. paragraphs (Point, Evidence, Explanation)."
    ]
    for pr in previews:
        p = tf.add_paragraph()
        p.text = f"👉 {pr}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_before = Pt(10)
    
    add_footer(slide)
    return slide


def build_slide_25_celebration(prs, total):
    """Slide 25: Celebration Slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, ROYAL_BLUE)
    
    # Gold border card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.5))
    card.fill.solid()
    card.fill.fore_color.rgb = ROYAL_BLUE
    card.line.color.rgb = GOLD
    card.line.width = Pt(3)
    
    tb = slide.shapes.add_textbox(Inches(1.8), Inches(1.3), Inches(9.733), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "🎉 EXCELLENT WORK TODAY! 🎉"
    p0.font.size = Pt(32)
    p0.font.bold = True
    p0.font.color.rgb = GOLD
    p0.alignment = PP_ALIGN.CENTER
    
    p1 = tf.add_paragraph()
    p1.text = "You have completed Lesson 1 of Key Stage 3 English."
    p1.font.size = Pt(18)
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(14)
    
    p_q = tf.add_paragraph()
    p_q.text = '"The more that you read, the more things you will know.\nThe more that you learn, the more places you\'ll go."'
    p_q.font.size = Pt(16)
    p_q.font.italic = True
    p_q.font.color.rgb = RGBColor(226, 232, 240)
    p_q.alignment = PP_ALIGN.CENTER
    p_q.space_before = Pt(24)
    
    p_auth = tf.add_paragraph()
    p_auth.text = "— Dr. Seuss"
    p_auth.font.size = Pt(14)
    p_auth.font.bold = True
    p_auth.font.color.rgb = GOLD
    p_auth.alignment = PP_ALIGN.CENTER
    p_auth.space_before = Pt(6)
    
    add_footer(slide, "BLUESKY EDUCATIONAL CONSULTS | Learning Beyond Borders")
    return slide


# ==========================================
# MAIN GENERATION
# ==========================================

def generate_presentation():
    """Generate the complete Year 7 English Lesson 1 presentation."""
    print("Generating Year 7 English Lesson 1 presentation...")
    
    prs = create_base_presentation()
    total_slides = 25
    
    # Build all 25 slides
    build_slide_1_cover(prs)
    build_slide_2_welcome(prs, total_slides)
    build_slide_3_icebreaker(prs, total_slides)
    build_slide_4_objectives(prs, total_slides)
    build_slide_5_success(prs, total_slides)
    build_slide_6_warmup(prs, total_slides)
    build_slide_7_what_is_english(prs, total_slides)
    build_slide_8_four_skills(prs, total_slides)
    build_slide_9_why_english(prs, total_slides)
    build_slide_10_vocabulary(prs, total_slides)
    build_slide_11_reading(prs, total_slides)
    build_slide_12_modelling(prs, total_slides)
    build_slide_13_guided(prs, total_slides)
    build_slide_14_comprehension(prs, total_slides)
    build_slide_15_vocab_activity(prs, total_slides)
    build_slide_16_think_pair_share(prs, total_slides)
    build_slide_17_independent(prs, total_slides)
    build_slide_18_creative(prs, total_slides)
    build_slide_19_mistakes(prs, total_slides)
    build_slide_20_real_life(prs, total_slides)
    build_slide_21_summary(prs, total_slides)
    build_slide_22_exit(prs, total_slides)
    build_slide_23_homework(prs, total_slides)
    build_slide_24_preview(prs, total_slides)
    build_slide_25_celebration(prs, total_slides)
    
    filename = "Year_7_English_Lesson_1_Premium.pptx"
    prs.save(filename)
    print(f"✅ SUCCESS! Presentation saved as '{filename}'")
    return filename


if __name__ == "__main__":
    generate_presentation()
